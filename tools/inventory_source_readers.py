"""Format-aware, source-preserving readers for inventory recovery."""

from __future__ import annotations

import csv
import io
import subprocess
from pathlib import Path
from typing import Any

import fitz
import openpyxl
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def read_source(path: Path, ocr_container: str) -> dict[str, Any]:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return read_pdf(path, ocr_container)
    if suffix == ".xlsx":
        return read_xlsx(path, ocr_container)
    if suffix == ".pptx":
        return read_pptx(path, ocr_container)
    raise ValueError(f"Unsupported inventory source: {suffix}")


def read_pdf(path: Path, ocr_container: str) -> dict[str, Any]:
    fragments: list[dict[str, Any]] = []
    tables: list[dict[str, Any]] = []
    assets: list[dict[str, Any]] = []
    with fitz.open(path) as document:
        for page_index, page in enumerate(document, start=1):
            page_text = 0
            blocks = page.get_text("dict", sort=True).get("blocks", [])
            for block_index, block in enumerate(blocks, start=1):
                if block.get("type") == 0:
                    text = "\n".join(
                        "".join(span.get("text", "")
                                for span in line.get("spans", []))
                        for line in block.get("lines", [])
                    ).strip()
                    if text:
                        page_text += len(text)
                        fragments.append(fragment(
                            f"pdf:page={page_index};block={block_index}",
                            "TEXT", text, block.get("bbox"), page_index,
                        ))
                elif block.get("type") == 1 and block.get("image"):
                    ocr = ocr_image(block["image"], ocr_container)
                    locator = f"pdf:page={page_index};image={block_index}"
                    assets.append(asset(locator, block, ocr, page_index))
                    if ocr["text"]:
                        fragments.append(fragment(
                            locator, "IMAGE_OCR", ocr["text"],
                            block.get("bbox"), page_index, ocr["confidence"],
                        ))
            for table_index, table in enumerate(
                    find_pdf_tables(page), start=1):
                rows = [[value or "" for value in row]
                        for row in table.extract()]
                tables.append({
                    "locator": f"pdf:page={page_index};table={table_index}",
                    "page": page_index,
                    "bbox": list(table.bbox),
                    "rows": rows,
                })
            if page_text < 20:
                pixmap = page.get_pixmap(
                    matrix=fitz.Matrix(2, 2), alpha=False)
                ocr = ocr_image(pixmap.tobytes("png"), ocr_container)
                if ocr["text"]:
                    fragments.append(fragment(
                        f"pdf:page={page_index};render-ocr=1",
                        "PAGE_OCR", ocr["text"],
                        [0, 0, page.rect.width, page.rect.height],
                        page_index, ocr["confidence"],
                    ))
    return result("PDF", fragments, tables, assets)


def find_pdf_tables(page: fitz.Page) -> list[Any]:
    try:
        return list(page.find_tables().tables)
    except (AttributeError, TypeError, ValueError):
        return []


def read_xlsx(path: Path, ocr_container: str) -> dict[str, Any]:
    formulas = openpyxl.load_workbook(
        path, data_only=False, read_only=False)
    cached = openpyxl.load_workbook(
        path, data_only=True, read_only=False)
    fragments: list[dict[str, Any]] = []
    tables: list[dict[str, Any]] = []
    assets: list[dict[str, Any]] = []
    try:
        for sheet_index, sheet in enumerate(
                formulas.worksheets, start=1):
            cached_sheet = cached.worksheets[sheet_index - 1]
            populated = [
                row for row in sheet.iter_rows()
                if any(cell.value is not None or
                       cached_sheet[cell.coordinate].value is not None
                       for cell in row)
            ]
            if populated:
                rows = []
                for row in populated:
                    values = []
                    for cell in row:
                        formula = (str(cell.value)
                                   if cell.data_type == "f" else None)
                        cached_value = cached_sheet[cell.coordinate].value
                        displayed = (cached_value
                                     if formula is not None and
                                     cached_value is not None
                                     else cell.value)
                        values.append({
                            "coordinate": cell.coordinate,
                            "value": (None if displayed is None
                                      else str(displayed)),
                            "formula": formula,
                            "cachedValue": (
                                None if cached_value is None
                                else str(cached_value)
                            ),
                            "numberFormat": cell.number_format,
                            "styleId": cell.style_id,
                        })
                    rows.append(values)
                tables.append({
                    "locator":
                        f"xlsx:sheet={sheet_index}:{sheet.title}",
                    "sheet": sheet.title,
                    "mergedRanges": [
                        str(item) for item in sheet.merged_cells.ranges
                    ],
                    "rows": rows,
                })
            for image_index, image in enumerate(
                    getattr(sheet, "_images", []), start=1):
                blob = image._data()
                anchor = getattr(
                    getattr(image, "anchor", None), "_from", None)
                cell = (f"{anchor.col + 1},{anchor.row + 1}"
                        if anchor else "unknown")
                locator = (
                    f"xlsx:sheet={sheet_index}:{sheet.title};"
                    f"image={image_index};anchor={cell}"
                )
                ocr = ocr_image(blob, ocr_container)
                assets.append({
                    "locator": locator,
                    "bytes": len(blob),
                    "ocrText": ocr["text"],
                    "ocrConfidence": ocr["confidence"],
                })
                if ocr["text"]:
                    fragments.append(fragment(
                        locator, "IMAGE_OCR", ocr["text"],
                        None, sheet_index, ocr["confidence"],
                    ))
    finally:
        formulas.close()
        cached.close()
    return result("XLSX", fragments, tables, assets)


def read_pptx(path: Path, ocr_container: str) -> dict[str, Any]:
    presentation = Presentation(path)
    fragments: list[dict[str, Any]] = []
    tables: list[dict[str, Any]] = []
    assets: list[dict[str, Any]] = []
    for slide_index, slide in enumerate(
            presentation.slides, start=1):
        read_pptx_shapes(
            slide.shapes, slide_index, "shape", ocr_container,
            fragments, tables, assets,
        )
    return result("PPTX", fragments, tables, assets)


def read_pptx_shapes(
    shapes: Any,
    slide_index: int,
    path_prefix: str,
    ocr_container: str,
    fragments: list[dict[str, Any]],
    tables: list[dict[str, Any]],
    assets: list[dict[str, Any]],
) -> None:
    for shape_index, shape in enumerate(shapes, start=1):
        locator = (
            f"pptx:slide={slide_index};"
            f"{path_prefix}={shape_index}"
        )
        bbox = [shape.left, shape.top, shape.width, shape.height]
        if (getattr(shape, "has_text_frame", False) and
                shape.text.strip()):
            fragments.append(fragment(
                locator, "TEXT", shape.text.strip(),
                bbox, slide_index,
            ))
        if getattr(shape, "has_table", False):
            tables.append({
                "locator": locator,
                "slide": slide_index,
                "bbox": bbox,
                "rows": [
                    [cell.text.strip() for cell in row.cells]
                    for row in shape.table.rows
                ],
            })
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            blob = shape.image.blob
            ocr = ocr_image(blob, ocr_container)
            assets.append({
                "locator": locator,
                "bytes": len(blob),
                "contentType": shape.image.content_type,
                "ocrText": ocr["text"],
                "ocrConfidence": ocr["confidence"],
            })
            if ocr["text"]:
                fragments.append(fragment(
                    locator, "IMAGE_OCR", ocr["text"],
                    bbox, slide_index, ocr["confidence"],
                ))
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            read_pptx_shapes(
                shape.shapes, slide_index,
                path_prefix + "." + str(shape_index) + ".shape",
                ocr_container, fragments, tables, assets,
            )


def ocr_image(blob: bytes, container: str) -> dict[str, Any]:
    if len(blob) < 1_024:
        return {"text": "", "confidence": None}
    completed = subprocess.run(
        [
            "docker", "exec", "-i", container,
            "tesseract", "stdin", "stdout", "-l", "eng", "tsv",
        ],
        input=blob,
        stdout=subprocess.PIPE,
        stderr=subprocess.DEVNULL,
        check=False,
        timeout=120,
    )
    if completed.returncode != 0:
        return {"text": "", "confidence": None}
    output = completed.stdout.decode(
        "utf-8", errors="replace")
    rows = list(csv.DictReader(
        io.StringIO(output), delimiter="\t"))
    words = [
        row.get("text", "").strip()
        for row in rows
        if row.get("text", "").strip()
    ]
    scores = []
    for row in rows:
        try:
            score = float(row.get("conf", "-1"))
        except (TypeError, ValueError):
            continue
        if score >= 0:
            scores.append(score / 100)
    return {
        "text": " ".join(words),
        "confidence": min(scores) if scores else None,
    }


def fragment(
    locator: str,
    kind: str,
    text: str,
    bbox: Any,
    ordinal: int,
    confidence: float | None = None,
) -> dict[str, Any]:
    return {
        "locator": locator,
        "kind": kind,
        "text": text,
        "bbox": list(bbox) if bbox else None,
        "ordinal": ordinal,
        "confidence": confidence,
    }


def asset(
    locator: str,
    block: dict[str, Any],
    ocr: dict[str, Any],
    page: int,
) -> dict[str, Any]:
    return {
        "locator": locator,
        "page": page,
        "bytes": len(block.get("image", b"")),
        "extension": block.get("ext"),
        "bbox": list(block.get("bbox", [])),
        "ocrText": ocr["text"],
        "ocrConfidence": ocr["confidence"],
    }


def result(
    format_name: str,
    fragments: list[dict[str, Any]],
    tables: list[dict[str, Any]],
    assets: list[dict[str, Any]],
) -> dict[str, Any]:
    return {
        "format": format_name,
        "fragments": fragments,
        "tables": tables,
        "assets": assets,
        "counts": {
            "fragments": len(fragments),
            "tables": len(tables),
            "assets": len(assets),
        },
    }
